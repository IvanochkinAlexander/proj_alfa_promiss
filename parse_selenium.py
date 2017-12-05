#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pymorphy2 import MorphAnalyzer
from nltk.tokenize import sent_tokenize
from nltk.tokenize import word_tokenize
from bs4 import BeautifulSoup as bs
import urllib
import re
import pandas as pd
import numpy as np
from selenium import webdriver
from bs4 import BeautifulSoup
import time
from collections import  Counter
import time
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.keys import Keys
import re
import datetime
from cStringIO import StringIO
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
import subprocess

date_time = datetime.datetime.now()

random_tags = ["panel",
              "news-text news-text-html",
              "layout__content content","sbrf-rich-wrapper",
              "content",
              "personalizable-container__cell col-xs-12 col-sm-6 col-md-3 personalizable-container__cell_occupied",
              "bp-widget-body",
              "card-info",
              "content-text",
              "mid_col",
              "wrapper",
              "wrapper__content",
              "section",
             "ui-trading-stock-about",
             "news-about__insurance",
             "news-about__desc  static-text",
             "bull",
             "columns__col  columns__col_content",
             "panel panel_news-about",
             "news-list__item  news-list__item_right-list  ",
             "columns__col panel columns__col_sidebar",
             "ui-collapse-banner__inner ui-collapse-banner__inner_children TCS-menu-float",
             "ui-action-menu",
             "ui-menu-second ui-menu-second_product",
             "ui-menu-second__layout",
             "ui-menu-second__wrapper ui-menu-second__wrapper_product",
             "ui-menu-second__container ui-layout__container ui-menu-second__container_product",
             "ui-scroll ui-menu-second__list ui-menu-second__list_product",
             "ui-scroll__window",
             "ui-header-primary__column ui-header-primary__column_toggle",
             "ui-header-primary__column ui-header-primary__column_menu",
             "ui-trading-news-list__content-announce",
             "ui-trading-security-news",
             "ui-trading-security-news__filte",
             "ui-context-filter ui-context-filter_filter-news",
             "application",
             "col-group",
             "main__inner main__inner_no-bg",
             "services-payment__item-text",
             "main_text",
             "b-container",
             "body-container",
             "descrip"
             ]

def convert(fname, pages=None):

    """Process pdf file"""

    if not pages:
        pagenums = set()
    else:
        pagenums = set(pages)

    output = StringIO()
    manager = PDFResourceManager()
    converter = TextConverter(manager, output, laparams=LAParams())
    interpreter = PDFPageInterpreter(manager, converter)

    infile = file(fname, 'rb')
    for page in PDFPage.get_pages(infile, pagenums):
        interpreter.process_page(page)
    infile.close()
    converter.close()
    text = output.getvalue()
    output.close
    return text

def make_unicode(input):

    """Convert string to unicode"""

    if type(input) != unicode:
        input =  input.decode('utf-8')
        return input
    else:
        return input


import docx
import docx2txt
from pptx import Presentation

def read_docx (file):

    """Process docx files"""

    text = docx2txt.process(file)
    return text

def read_pptx (file):

    """Process powerpoint files"""

    prs = Presentation(file)

    text_runs = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs+=run.text
                    text_runs+=' '
    return text_runs


def cut_data (links_df):

    """Reduces number of samples by 150 and cuts long domain names"""

    links_df['length'] = links_df[0].apply(lambda x: x.count('/'))
    links_df['count'] = links_df.groupby('domain')[0].transform('count')
    links_df = links_df.sort_values(['domain', 'length'])
    final_df = pd.DataFrame()
    for i in links_df['domain'].unique():
        links_df_df = links_df[links_df['domain']==i]
        links_df_df= links_df_df[:300]
        final_df= pd.concat([final_df, links_df_df])
    final_df = final_df.reset_index(drop=True)
    final_df = final_df[final_df['length']<=6]
    return final_df

def parse_links (links_alfa, random_tags):

    """Parse all links"""

    long_text = []
    values = []
    for i, k in enumerate(pd.DataFrame(links_alfa)[0].unique()):
        try:
            print (k)
            if k.split('.')[-1] == 'pdf':
                name = k.split('/')[-1]
                subprocess.check_output('!wget -N {k} -P ../output/promiss/pdfs') ### -N проверят версию файла, если она новее, то файл перезаписывается
                text_from_pdf = convert('pdfs/{}'.format(name))
                print (i)
                long_text.append(i)
                long_text.append(text_from_pdf)
            elif k.split('.')[-1] == 'docx':
                name = k.split('/')[-1]
                subprocess.check_output('!wget -N {k} -P ../output/promiss/pdfs') ### -N проверят версию файла, если она новее, то файл перезаписывается
                text_from_pdf = read_docx('pdfs/{}'.format(name))
                print (i)
                long_text.append(i)
                long_text.append(text_from_pdf)
            elif k.split('.')[-1] == 'pptx':
                name = k.split('/')[-1]
                subprocess.check_output('!wget -N {k} -P ../output/promiss/pdfs') ### -N проверят версию файла, если она новее, то файл перезаписывается
                text_from_pdf = read_pptx('pdfs/{}'.format(name))
                print (i)
                long_text.append(i)
                long_text.append(text_from_pdf)
            else:
                YOUR_PAGE_URL = k

                browser.get(YOUR_PAGE_URL)
                page = browser.page_source
                soup = BeautifulSoup(page, "lxml")
                for value in soup.find_all('div', class_=random_tags):
                    long_text.append(i)
                    long_text.append(value.text)

        except:
            pass
    return long_text



def post_processing (long_text):

    """Convert parsed links to dataframe"""

    text_down = pd.DataFrame(long_text)
    text_down['shifted'] = text_down[0].shift(1)
    text_down['type'] = text_down[0].apply(lambda x:type(x))
    text_down = text_down[(text_down['type']==unicode) | (text_down['type']==str)]
    text_down=text_down.rename(columns={0:'text'})
    links_df = pd.DataFrame(pd.DataFrame(links_alfa)[0].unique()).reset_index()
    merged_text = pd.merge(text_down, links_df, left_on='shifted', right_on='index', how='left')
    merged_text_cut = merged_text[['text', 0]]
    merged_text_cut.columns=('full_text', 'link')
    merged_text_cut = merged_text_cut.drop_duplicates(subset='full_text', keep='first')
    merged_text_cut['full_text'] = merged_text_cut['full_text'].apply(make_unicode)
    merged_text_cut.columns = ('text', 'url')
    merged_text_cut[['url', 'text']].to_excel('../output/promiss/selenium_parsed.xlsx', index=False)
    print (merged_text_cut.shape[0])


links_alfa = pd.read_excel('../output/promiss/to_parse_selenium.xlsx')
links_alfa = cut_data(links_alfa)
links_alfa.to_excel('../output/promiss/to_parse_selenium.xlsx', index=False)
print ('новая размерность составляет {}'.format(links_alfa.shape[0]))
browser = webdriver.PhantomJS(executable_path='../../../usr/local/bin/phantomjs')
browser.set_page_load_timeout(15)
long_text = parse_links(links_alfa, random_tags)
post_processing (long_text)
